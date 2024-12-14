import openai
import pptx
from pptx.util import Inches, Pt
import os
import base64
import streamlit as st
from dotenv import load_dotenv

# Load OpenAI API Key
load_dotenv()
openai.api_key = os.getenv('OPENAI_API_KEY')

# Define custom formatting options
TITLE_FONT_SIZE = Pt(25)
SLIDE_FONT_SIZE = Pt(20)
MAX_TEXT_LENGTH = 400  # Maximum length of text before reducing font size

# Function to Generate the slide titles
def generate_slide_titles(topic):
    try:
        messages = [
            {"role": "system", "content": "You are an AI assistant that helps create presentations."},
            {"role": "user", "content": f"Generate 10 slide titles for the topic '{topic}'."}
        ]
        response = openai.ChatCompletion.create(
            model="gpt-4",
            messages=messages,
            max_tokens=100,
        )
        return response['choices'][0]['message']['content'].split("\n")
    except Exception as e:
        st.error(f"Error generating slide titles: {e}")
        return []

# Function to Generate the slide content for each title
def generate_slide_content(slide_title):
    try:
        messages = [
            {"role": "system", "content": "You are an AI assistant that helps create presentations."},
            {"role": "user", "content": f"Generate content for the slide: '{slide_title}'."}
        ]
        response = openai.ChatCompletion.create(
            model="gpt-4",
            messages=messages,
            max_tokens=150,  # Adjust as needed based on the desired content length
        )
        return response['choices'][0]['message']['content']
    except Exception as e:
        st.error(f"Error generating slide content for '{slide_title}': {e}")
        return ""

# Function to adjust font size dynamically based on text length
def adjust_font_size(text_frame, text, max_chars=MAX_TEXT_LENGTH):
    """Adjust font size dynamically based on text length."""
    font_size = SLIDE_FONT_SIZE

    if len(text) > max_chars:
        font_size = Pt(12)  # Smaller font size if text is too long
    elif len(text) > max_chars // 2:
        font_size = Pt(14)  # Medium font size if text is moderately long

    for paragraph in text_frame.paragraphs:
        paragraph.font.size = font_size

# Function to Generate the presentation
def create_presentation(topic, slide_titles, slide_contents):
    # Create the 'generated_ppt' directory if it doesn't exist
    output_directory = "generated_ppt"
    os.makedirs(output_directory, exist_ok=True)

    prs = pptx.Presentation()
    slide_layout = prs.slide_layouts[1]

    # Title slide
    title_slide = prs.slides.add_slide(prs.slide_layouts[0])
    title_slide.shapes.title.text = topic
    title_text_frame = title_slide.shapes.title.text_frame
    adjust_font_size(title_text_frame, topic)  # Adjust title font size

    # Content slides
    for slide_title, slide_content in zip(slide_titles, slide_contents):
        slide = prs.slides.add_slide(slide_layout)
        slide.shapes.title.text = slide_title
        slide_content_shape = slide.shapes.placeholders[1]
        slide_content_shape.text = slide_content

        # Adjust font size for slide content
        adjust_font_size(slide_content_shape.text_frame, slide_content)

    # Save the presentation in the created directory
    presentation_path = os.path.join(output_directory, f"{topic}_presentation.pptx")
    prs.save(presentation_path)

    return presentation_path  # Return the full path to be used for download

# Streamlit app
def main():
    st.title("Text to PowerPoint :computer:")

    topic = st.text_input("Enter the topic for your presentation:")
    generate_button = st.button("Generate Presentation")

    if generate_button and topic:
        st.info("Generating presentation... Please wait.")
        slide_titles = generate_slide_titles(topic)
        filtered_slide_titles = [item for item in slide_titles if item.strip() != '']
        print("Slide Title: ", filtered_slide_titles)
        slide_contents = [generate_slide_content(title) for title in filtered_slide_titles]
        print("Slide Contents: ", slide_contents)
        presentation_path = create_presentation(topic, filtered_slide_titles, slide_contents)
        print("Presentation generated successfully!")

        st.success("Presentation generated successfully!")
        st.download_button(
            label="Download the PowerPoint Presentation",
            data=open(presentation_path, "rb").read(),
            file_name=f"{topic}_presentation.pptx",
            mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
        )

if __name__ == "__main__":
    main()
